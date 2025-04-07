from fastapi import FastAPI, File, UploadFile
import subprocess
import whisper
import multiprocessing
import os
from vosk import Model, KaldiRecognizer
import wave
import json
from docx import Document
from pptx import Presentation
from PyPDF2 import PdfReader

app = FastAPI()

# Load Whisper Model
whisper_model = whisper.load_model("base")

# Load Vosk Model
vosk_model_path = "model"
if not os.path.exists(vosk_model_path):
    raise Exception("Vosk model not found. Download from https://alphacephei.com/vosk/models")
vosk_model = Model(vosk_model_path)

ALLOWED_EXTENSIONS = ('.mp4', '.mkv', '.avi', '.docx', '.pptx', '.pdf')

def extract_audio(video_path, audio_path):
    command = ["ffmpeg", "-i", video_path, "-ac", "1", "-ar", "16000", audio_path, "-y"]
    process = subprocess.run(command, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
    if process.returncode != 0:
        raise Exception(f"FFmpeg failed: {process.stderr.decode()}")

def transcribe_whisper(audio_path):
    result = whisper_model.transcribe(audio_path)
    return result["text"]

def transcribe_vosk(audio_path):
    wf = wave.open(audio_path, "rb")
    rec = KaldiRecognizer(vosk_model, wf.getframerate())
    transcript = ""
    while True:
        data = wf.readframes(4000)
        if len(data) == 0:
            break
        if rec.AcceptWaveform(data):
            transcript += json.loads(rec.Result())["text"]
    return transcript

def extract_text_from_docx(file_path):
    doc = Document(file_path)
    text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
    return text

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text_from_pdf(file_path):
    with open(file_path, "rb") as f:
        reader = PdfReader(f)
        text = ""
        for page in reader.pages:
            text += page.extract_text() + "\n"
    return text

@app.post("/transcribe/")
async def transcribe_file(file: UploadFile = File(...), method: str = "whisper"):
    if not file.filename.endswith(ALLOWED_EXTENSIONS):
        return {"error": "Unsupported file type"}

    file_path = f"temp_{file.filename}"

    try:
        with open(file_path, "wb") as f:
            content = await file.read()
            f.write(content)

        if file.filename.endswith(('.mp4', '.mkv', '.avi')):
            audio_path = os.path.splitext(file_path)[0] + ".wav"
            extract_audio(file_path, audio_path)
            text = transcribe_whisper(audio_path) if method == "whisper" else transcribe_vosk(audio_path)
        elif file.filename.endswith('.docx'):
            text = extract_text_from_docx(file_path)
        elif file.filename.endswith('.pptx'):
            text = extract_text_from_pptx(file_path)
        elif file.filename.endswith('.pdf'):
            text = extract_text_from_pdf(file_path)
        else:
            return {"error": "Unsupported file type"}

    except Exception as e:
        return {"error": str(e)}

    finally:
        if os.path.exists(file_path):
            os.remove(file_path)
        if 'audio_path' in locals() and os.path.exists(audio_path):
            os.remove(audio_path)

    return {"transcription": text}

@app.on_event("shutdown")
def cleanup():
    print("Cleaning up resources...")
    global whisper_model, vosk_model
    whisper_model = None
    vosk_model = None
    multiprocessing.active_children()
    for child in multiprocessing.active_children():
        child.terminate()
    os.system("sync")